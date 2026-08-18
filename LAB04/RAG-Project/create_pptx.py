import os
import sys

def install_and_import():
    try:
        import pptx
    except ImportError:
        print("Installing python-pptx...")
        import subprocess
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx

install_and_import()

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize presentation with Widescreen aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# -------------------------------------------------------------
# Slide 1: RAG System Workflow Diagram
# -------------------------------------------------------------
slide1 = prs.slides.add_slide(prs.slide_layouts[5]) # Blank slide with title
title1 = slide1.shapes.title
title1.text = "RAG System Workflow"
title1.text_frame.paragraphs[0].font.size = Pt(36)
title1.text_frame.paragraphs[0].font.bold = True

# Helper to draw colored boxes
def create_box(slide, text, subtext, left, top, width, height, bg_r, bg_g, bg_b, line_r, line_g, line_b, dashed=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(bg_r, bg_g, bg_b)
    shape.line.color.rgb = RGBColor(line_r, line_g, line_b)
    shape.line.width = Pt(2.5)
    
    if dashed:
        from pptx.enum.dml import MSO_LINE
        shape.line.dash_style = MSO_LINE.DASH
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.alignment = PP_ALIGN.CENTER
    
    if subtext:
        p2 = tf.add_paragraph()
        p2.text = subtext
        p2.font.bold = False
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(90, 90, 90)
        p2.alignment = PP_ALIGN.CENTER

    return shape

# Helper to draw arrows
def add_arrow(slide, start_x, start_y, end_x, end_y):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, start_x, start_y, end_x - start_x, Inches(0.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
    shape.line.color.rgb = RGBColor(200, 200, 200)
    return shape

# Helper to draw section backgrounds
def add_section_bg(slide, title, top, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), top, Inches(12.7), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(220, 220, 220)
    shape.line.width = Pt(1.5)
    
    txBox = slide.shapes.add_textbox(Inches(0.4), top - Inches(0.2), Inches(5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(80, 80, 80)

# Colors matching the requested visual style
blue_bg, blue_line = (239, 246, 255), (96, 165, 250)
orange_bg, orange_line = (255, 247, 237), (249, 115, 22)
green_bg, green_line = (240, 253, 244), (74, 222, 128)
purple_bg, purple_line = (250, 245, 255), (192, 132, 252)

# --- SECTION 1: BUILD ---
add_section_bg(slide1, "1 BUILD — Prepare dataset", Inches(1.4), Inches(1.5))
create_box(slide1, "Knowledge Base", "data/diabetes_q_a.txt", Inches(0.5), Inches(1.7), Inches(1.6), Inches(0.8), *blue_bg, *blue_line)
add_arrow(slide1, Inches(2.1), Inches(2.0), Inches(2.5), Inches(2.0))
create_box(slide1, "document_loader", "Parse files from Q&A", Inches(2.6), Inches(1.7), Inches(1.7), Inches(0.8), *blue_bg, *blue_line)
add_arrow(slide1, Inches(4.3), Inches(2.0), Inches(4.7), Inches(2.0))
create_box(slide1, "text_splitter", "Split text into chunks", Inches(4.8), Inches(1.7), Inches(1.7), Inches(0.8), *blue_bg, *blue_line)
add_arrow(slide1, Inches(6.5), Inches(2.0), Inches(6.9), Inches(2.0))
create_box(slide1, "embedding_model", "MiniLM multilingual", Inches(7.0), Inches(1.7), Inches(1.7), Inches(0.8), *blue_bg, *blue_line)

add_arrow(slide1, Inches(8.7), Inches(1.7), Inches(9.1), Inches(1.7))
add_arrow(slide1, Inches(8.7), Inches(2.3), Inches(9.1), Inches(2.3))

create_box(slide1, "vector_store", "FAISS index", Inches(9.2), Inches(1.45), Inches(1.4), Inches(0.6), *blue_bg, *blue_line)
create_box(slide1, "build_bm25", "hybrid_retriever", Inches(9.2), Inches(2.15), Inches(1.4), Inches(0.6), *blue_bg, *blue_line)

add_arrow(slide1, Inches(10.6), Inches(2.0), Inches(11.0), Inches(2.0))
create_box(slide1, "vector_db/", "document.index\nbm25_index.pkl", Inches(11.1), Inches(1.7), Inches(1.7), Inches(0.8), *purple_bg, *purple_line)


# --- SECTION 2: QUERY ---
add_section_bg(slide1, "2 QUERY — Answer questions (Input → Retrieval → Context → LLM → Output)", Inches(3.3), Inches(1.8))
create_box(slide1, "[INPUT]\nUser Question", "User asks a question", Inches(0.5), Inches(3.6), Inches(1.5), Inches(0.9), *green_bg, *green_line)
add_arrow(slide1, Inches(2.0), Inches(4.0), Inches(2.4), Inches(4.0))
create_box(slide1, "query_transform", "Rewrite, multi, HyDE", Inches(2.5), Inches(3.6), Inches(1.7), Inches(0.9), *green_bg, *green_line, dashed=True)
add_arrow(slide1, Inches(4.2), Inches(4.0), Inches(4.6), Inches(4.0))
create_box(slide1, "[RETRIEVAL]\nhybrid_retriever", "BM25 + dense → RRF", Inches(4.7), Inches(3.6), Inches(1.9), Inches(0.9), *orange_bg, *orange_line)
add_arrow(slide1, Inches(6.6), Inches(4.0), Inches(7.0), Inches(4.0))
create_box(slide1, "[CONTEXT]\nprompt_templates", "Combine chunks &\nConversation Memory", Inches(7.1), Inches(3.6), Inches(2.0), Inches(0.9), *orange_bg, *orange_line)
add_arrow(slide1, Inches(9.1), Inches(4.0), Inches(9.5), Inches(4.0))
create_box(slide1, "[LLM]\ngenerator", "OpenAI / Ollama / Gemini", Inches(9.6), Inches(3.6), Inches(1.8), Inches(0.9), *orange_bg, *orange_line)
add_arrow(slide1, Inches(11.4), Inches(4.0), Inches(11.7), Inches(4.0))
create_box(slide1, "[OUTPUT]\nAnswer", "Response + citations", Inches(11.8), Inches(3.6), Inches(1.0), Inches(0.9), *green_bg, *green_line)


# --- SECTION 3: EVALUATION ---
add_section_bg(slide1, "3 EVALUATION — Measure and improve", Inches(5.5), Inches(1.5))
create_box(slide1, "golden_set.json", "60 questions", Inches(0.5), Inches(5.8), Inches(1.6), Inches(0.8), *purple_bg, *purple_line)
add_arrow(slide1, Inches(2.1), Inches(6.2), Inches(2.5), Inches(6.2))
create_box(slide1, "eval_retrieval.py", "Hit@k, MRR, nDCG", Inches(2.6), Inches(5.8), Inches(1.8), Inches(0.8), *blue_bg, *blue_line)
add_arrow(slide1, Inches(4.4), Inches(6.2), Inches(4.8), Inches(6.2))
create_box(slide1, "Retrieval Report", "outputs/eval_retrieval.json", Inches(4.9), Inches(5.8), Inches(2.0), Inches(0.8), *purple_bg, *purple_line)

create_box(slide1, "eval_generation.py", "faithfulness, correctness", Inches(7.4), Inches(5.8), Inches(2.0), Inches(0.8), *blue_bg, *blue_line)
add_arrow(slide1, Inches(9.4), Inches(6.2), Inches(9.8), Inches(6.2))
create_box(slide1, "Generation Report", "outputs/eval_generation.json", Inches(9.9), Inches(5.8), Inches(2.1), Inches(0.8), *purple_bg, *purple_line)

# -------------------------------------------------------------
# Slide 2: Workflow Explanation
# -------------------------------------------------------------
slide2 = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
title2 = slide2.shapes.title
title2.text = "RAG System Workflow Explanation"

content = slide2.placeholders[1]
tf2 = content.text_frame
tf2.text = "1. Input"
p1 = tf2.add_paragraph()
p1.text = "User submits a question via the terminal (main.py)."
p1.level = 1
p1_1 = tf2.add_paragraph()
p1_1.text = "The query_transform module replaces slang and optionally rewrites the query using an LLM."
p1_1.level = 1

p2 = tf2.add_paragraph()
p2.text = "2. Retrieval"
p2.level = 0
p2_1 = tf2.add_paragraph()
p2_1.text = "hybrid_retriever searches the knowledge base using FAISS (dense vectors) and BM25 (keyword matching)."
p2_1.level = 1
p2_2 = tf2.add_paragraph()
p2_2.text = "Results are fused together via Reciprocal Rank Fusion (RRF) to isolate the top 3 relevant chunks."
p2_2.level = 1

p3 = tf2.add_paragraph()
p3.text = "3. Context"
p3.level = 0
p3_1 = tf2.add_paragraph()
p3_1.text = "prompt_templates structures the retrieved chunks into a numbered context block (e.g., [1], [2])."
p3_1.level = 1
p3_2 = tf2.add_paragraph()
p3_2.text = "The context is seamlessly merged with the user's question and conversation memory."
p3_2.level = 1

p4 = tf2.add_paragraph()
p4.text = "4. LLM"
p4.level = 0
p4_1 = tf2.add_paragraph()
p4_1.text = "The generator transmits the context-enriched prompt to the LLM via the OpenAI python client."
p4_1.level = 1
p4_2 = tf2.add_paragraph()
p4_2.text = "The LLM is strictly instructed to answer using only the provided context and inject source citations."
p4_2.level = 1

p5 = tf2.add_paragraph()
p5.text = "5. Output"
p5.level = 0
p5_1 = tf2.add_paragraph()
p5_1.text = "The generated answer is parsed, and a mandatory medical disclaimer is appended."
p5_1.level = 1
p5_2 = tf2.add_paragraph()
p5_2.text = "The final response is displayed to the user and recorded into the conversation memory loop."
p5_2.level = 1

# Ensure text fits and isn't overcrowded
tf2.fit_text(font_family="Calibri", max_size=20)

prs.save("RAG_System_Workflow.pptx")
print("Successfully generated RAG_System_Workflow.pptx")
